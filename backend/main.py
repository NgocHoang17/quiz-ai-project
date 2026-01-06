from fastapi import FastAPI, Depends, HTTPException, status, File, UploadFile, Form
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session, joinedload 
from sqlalchemy import func, desc 
from typing import List, Optional
from datetime import timedelta
import io
import fitz
import docx
from pptx import Presentation
from fastapi.responses import StreamingResponse 
from docx.shared import Pt 
from docx.enum.text import WD_ALIGN_PARAGRAPH 
from datetime import datetime, timedelta, date 

from . import models, schemas, security
from .database import engine, SessionLocal
import google.generativeai as genai
from pydantic import BaseModel, Field
from dotenv import load_dotenv
import os
import json

#  TẢI BIẾN MÔI TRƯỜNG & CẤU HÌNH AI 
dotenv_path = os.path.join(os.path.dirname(__file__), '..', '.env')
load_dotenv(dotenv_path=dotenv_path)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print(" LỖI: Không tìm thấy GEMINI_API_KEY")
else:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        print("✅ Đã cấu hình thành công Google AI.")
    except Exception as e:
        print(f"❌ Lỗi khi cấu hình Google AI: {e}")

#  TẠO BẢNG TRONG DATABASE 
models.Base.metadata.create_all(bind=engine)

app = FastAPI()

#  CẤU HÌNH CORS 
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

#  HÀM PHỤ THUỘC (DEPENDENCIES) 
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="login")

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Không thể xác thực",
        headers={"WWW-Authenticate": "Bearer"},
    )
    token_data = security.verify_token(token, credentials_exception)
    user = db.query(models.User).filter(models.User.email == token_data.email).first()
    if user is None:
        raise credentials_exception
    return user

#  API XÁC THỰC (ĐĂNG KÝ & ĐĂNG NHẬP) 
@app.post("/register", response_model=schemas.UserOut)
def create_user(user: schemas.UserCreate, db: Session = Depends(get_db)):
    db_user = db.query(models.User).filter(models.User.email == user.email).first()
    if db_user:
        raise HTTPException(status_code=400, detail="Email đã được đăng ký")
    
    hashed_password = security.get_password_hash(user.password)
    new_user = models.User(email=user.email, hashed_password=hashed_password)
    
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    return new_user

@app.post("/login", response_model=schemas.Token)
def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(models.User).filter(models.User.email == form_data.username).first()
    
    if not user or not security.verify_password(form_data.password, user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Email hoặc mật khẩu không chính xác",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    access_token_expires = timedelta(minutes=security.ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = security.create_access_token(
        data={"sub": user.email}, expires_delta=access_token_expires
    )
    
    return {"access_token": access_token, "token_type": "bearer"}


#  HÀM LÕI AI 
def _generate_quiz_from_text(text: str, num_questions: int, quiz_type: str):
    quiz_type_instructions = {
        "mcq": "Trắc nghiệm 4 lựa chọn (A, B, C, D).",
        "fill_in_blank": "Trắc nghiệm dạng điền vào chỗ trống. Câu hỏi phải có một đến hai dấu '____' ở các vị trí quan trong trong câu trích dẫn để điền.",
        "exercise": "Câu hỏi dạng bài tập vận dụng hoặc giải quyết vấn đề dựa trên nội dung.(không cung cấp kí hiệu cho công thức và không giải thích dài dòng)",
        "mixed": "Hỗn hợp nhiều dạng câu hỏi (trắc nghiệm, điền khuyết, bài tập nếu có)."
    }
    instruction = quiz_type_instructions.get(quiz_type, quiz_type_instructions["mcq"])

    try:
        model = genai.GenerativeModel("gemini-2.5-flash") 
        
        prompt = f"""
        Dựa vào đoạn văn bản sau đây:
        "{text}"

       Hãy thực hiện các yêu cầu sau:
        1. Tạo chính xác {num_questions} câu hỏi loại: {instruction}
        2. Với mỗi câu hỏi, hãy viết một lời "giai_thich" ngắn gọn (tại sao đáp án đó đúng).
        3. Với mỗi câu hỏi, hãy tìm "trich_dan" là CÂU NGUYÊN VĂN trong văn bản chứa thông tin trả lời.
       Chú ý: Các câu hỏi không dùng các từ theo văn bản, theo đoạn văn, theo đoạn trích.

        QUY TẮC ĐỊNH DẠNG JSON:
        - Trả về CHUỖI JSON hợp lệ, KHÔNG thêm văn bản giải thích nào bên ngoài.
        - MỖI câu hỏi phải có 4 lựa chọn (A, B, C, D) và 1 đáp án đúng (dù là dạng điền khuyết hay bài tập).
        - Định dạng JSON phải là một danh sách (list) các đối tượng:
        [
          {{
            "cau_hoi": "Câu hỏi...?",
            "lua_chon": {{ "A": "...", "B": "...", "C": "...", "D": "..." }},
            "dap_an": "A",
            "giai_thich": "Giải thích ngắn gọn...",
            "trich_dan": "Trích dẫn nguyên văn từ tài liệu..."
          }}
        ]
        """
        response = model.generate_content(prompt)
        quiz_json_string = response.text.strip().replace("```json", "").replace("```", "").strip()
        
        try:
            quiz_data = json.loads(quiz_json_string)
            return quiz_data, None 
        except json.JSONDecodeError:
            print(f"DEBUG: AI trả về JSON không hợp lệ: {quiz_json_string}")
            return None, "AI trả về định dạng JSON không hợp lệ."

    except Exception as e:
        print(f" Lỗi khi gọi Gemini API: {e}")
        return None, f"Lỗi khi gọi Gemini API: {str(e)}"

#  API QUIZ (TẠO VÀ LƯU) 

class QuizRequest(BaseModel):
    text: str
    num_questions: int = Field(5, gt=0, le=25)
    quiz_type: str = "mcq"

@app.post("/generate-quiz") 
def generate_quiz_from_text(request: QuizRequest):
    input_text = request.text.strip()
    if not input_text:
        return {"error": "Vui lòng nhập nội dung văn bản."}

    quiz_data, error = _generate_quiz_from_text(
        request.text, request.num_questions, request.quiz_type
    )
    
    if error:
        return {"error": error}
    return {"quiz_data": quiz_data}

@app.post("/upload-quiz-file")
async def generate_quiz_from_file(
    file: UploadFile = File(...),
    num_questions: int = Form(5),
    quiz_type: str = Form("mcq")
):
    if not file:
        return {"error": "Vui lòng tải lên một file."}
        
    extracted_text = ""
    file_bytes = await file.read()
    file_extension = os.path.splitext(file.filename)[1].lower()
    
    try:
        if file_extension == ".txt":
            extracted_text = file_bytes.decode("utf-8")
        elif file_extension == ".pdf":
            with fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf") as doc:
                for page in doc:
                    extracted_text += page.get_text()
        elif file_extension == ".docx":
            doc_stream = io.BytesIO(file_bytes)
            doc = docx.Document(doc_stream)
            all_text = [p.text for p in doc.paragraphs]
            extracted_text = "\n".join(all_text)
        elif file_extension == ".pptx":
            ppt_stream = io.BytesIO(file_bytes)
            prs = Presentation(ppt_stream)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            extracted_text = "\n".join(all_text)
        else:
            return {"error": "Định dạng file không được hỗ trợ. Chỉ chấp nhận .txt, .pdf, .docx, .pptx."}
            
    except Exception as e:
        print(f"Lỗi khi đọc file {file.filename}: {e}")
        return {"error": f"Không thể đọc file. Lỗi: {str(e)}"}

    if not extracted_text.strip():
        return {"error": "File không có nội dung hoặc không thể trích xuất văn bản."}
    
    quiz_data, error = _generate_quiz_from_text(
        extracted_text, num_questions, quiz_type
    )
    
    if error:
        return {"error": error}
    return {"quiz_data": quiz_data}


#  API QUẢN LÝ QUIZ (LƯU, XEM, SỬA, XÓA) 

@app.post("/save-quiz", response_model=schemas.QuizOut)
def save_quiz(
    quiz_to_save: schemas.QuizCreate, 
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    try:
        new_quiz = models.Quiz(
            title=quiz_to_save.title,
            owner_id=current_user.id,
            folder_id=quiz_to_save.folder_id,
            quiz_type=quiz_to_save.quiz_type
        )
        db.add(new_quiz)
        db.commit()
        db.refresh(new_quiz)
        
        for q in quiz_to_save.questions:
            new_question = models.Question(
                question_text=q.question_text,
                choice_a=q.choice_a,
                choice_b=q.choice_b,
                choice_c=q.choice_c,
                choice_d=q.choice_d,
                correct_answer=q.correct_answer,
                explanation=q.explanation,
                citation=q.citation,
                quiz_id=new_quiz.id
            )
            db.add(new_question)
            
        db.commit() 
        db.refresh(new_quiz)
        
        return new_quiz
        
    except Exception as e:
        db.rollback() 
        print(f"Lỗi khi lưu quiz: {e}")
        raise HTTPException(status_code=500, detail=f"Lỗi máy chủ nội bộ: {e}")

@app.get("/my-quizzes", response_model=List[schemas.QuizOut])
def get_user_quizzes(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    """
    Trả về danh sách quiz VÀ các câu hỏi của chúng.
    """
    user_with_quizzes = db.query(models.User).options(
        joinedload(models.User.quizzes) # Tải các quiz
            .joinedload(models.Quiz.questions) # Tải các câu hỏi cho từng quiz
    ).filter(models.User.id == current_user.id).first()
    
    if not user_with_quizzes:
        return []
        
    return user_with_quizzes.quizzes

@app.delete("/quizzes/{quiz_id}")
def delete_quiz(
    quiz_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy bộ quiz")
    if quiz.owner_id != current_user.id:
        raise HTTPException(status_code=403, detail="Không có quyền xóa bộ quiz này")
    db.delete(quiz)
    db.commit()
    return {"message": "Đã xóa quiz thành công"}

@app.put("/quizzes/{quiz_id}", response_model=schemas.QuizOut)
def update_quiz_title(
    quiz_id: int,
    quiz_update: schemas.QuizBase,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy bộ quiz")
    if quiz.owner_id != current_user.id:
        raise HTTPException(status_code=403, detail="Không có quyền sửa bộ quiz này")
    
    quiz.title = quiz_update.title
    db.commit()
    db.refresh(quiz)
    return quiz

@app.get("/quizzes/{quiz_id}", response_model=schemas.QuizOut)
def get_quiz_details(
    quiz_id: int, 
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user) 
):
    """
    Lấy chi tiết một bộ quiz VÀ các câu hỏi của nó.
    """
    quiz = db.query(models.Quiz).options(
        joinedload(models.Quiz.questions) # Tải kèm câu hỏi
    ).filter(models.Quiz.id == quiz_id).first()
    
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy bộ quiz")
    
    if quiz.owner_id != current_user.id:
         raise HTTPException(status_code=403, detail="Bạn không có quyền truy cập Quiz này")

    return quiz


#  API NHẬN KẾT QUẢ VÀ LƯU LỊCH SỬ 
# Định nghĩa Model cho item trong list gửi lên (để tránh trùng tên với models.QuizResult)
class QuizSubmissionItem(BaseModel):
    question_id: int
    is_correct: bool

@app.post("/submit-quiz/{quiz_id}")
def submit_quiz_results(
    quiz_id: int, 
    results: List[QuizSubmissionItem], 
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    """
    1. Cập nhật UserQuestionStats (để biết câu nào hay sai).
    2. Lưu QuizResult (Lịch sử làm bài, điểm số).
    """
    correct_count = 0
    total_questions = len(results)

    try:
        # 1. Cập nhật thống kê từng câu 
        for result in results:
            if result.is_correct: 
                correct_count += 1
            
            stat = db.query(models.UserQuestionStats).filter(
                models.UserQuestionStats.user_id == current_user.id,
                models.UserQuestionStats.question_id == result.question_id
            ).first()
            
            if not stat:
                stat = models.UserQuestionStats(
                    user_id=current_user.id,
                    question_id=result.question_id,
                    correct_attempts=0,
                    incorrect_attempts=0
                )
                db.add(stat)
            
            if result.is_correct:
                stat.correct_attempts += 1
            else:
                stat.incorrect_attempts += 1
            
            stat.last_attempted_at = datetime.utcnow()

        # 2.  LƯU LỊCH SỬ LÀM BÀI 
        # Tính điểm thang 10
        score_10 = round((correct_count / total_questions) * 10) if total_questions > 0 else 0
        
        history_entry = models.QuizResult(
            user_id=current_user.id,
            quiz_id=quiz_id,
            score=score_10,
            total_questions=total_questions,
            completed_at=datetime.utcnow()
        )
        db.add(history_entry)
        
        db.commit()
    except Exception as e:
        db.rollback() 
        print(f"Lỗi khi submit quiz: {e}")
        raise HTTPException(status_code=500, detail="Lỗi máy chủ khi lưu kết quả.")

    return {"message": "Đã ghi nhận kết quả và lưu lịch sử."}

#  API STATS CHO CÂU HỎI HAY SAI (Dùng cho chế độ luyện tập) 
class QuestionStatOut(BaseModel):
    question_id: int
    is_frequently_wrong: bool

@app.get("/my-question-stats/{quiz_id}", response_model=List[QuestionStatOut])
def get_question_stats_for_quiz(
    quiz_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    question_ids = db.query(models.Question.id).filter(models.Question.quiz_id == quiz_id).all()
    q_ids = [id_tuple[0] for id_tuple in question_ids]
    stats = db.query(models.UserQuestionStats).filter(
        models.UserQuestionStats.user_id == current_user.id,
        models.UserQuestionStats.question_id.in_(q_ids)
    ).all()
    result_list = []
    for s in stats:
        is_wrong = (s.incorrect_attempts > 0) and (s.incorrect_attempts >= s.correct_attempts)
        result_list.append(
            QuestionStatOut(
                question_id=s.question_id,
                is_frequently_wrong=is_wrong
            )
        )
    return result_list

#  API DASHBOARD STATS  
@app.get("/dashboard-stats")
def get_dashboard_stats(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    # 1. Tổng số Quiz
    total_quizzes = db.query(models.Quiz).filter(models.Quiz.owner_id == current_user.id).count()
    
    # 2.  Tổng số lần làm bài (Thay cho tổng câu hỏi đã ôn)
    total_completed = db.query(models.QuizResult).filter(models.QuizResult.user_id == current_user.id).count()
    
    # 3.  Tổng số Folder
    total_folders = db.query(models.Folder).filter(models.Folder.user_id == current_user.id).count()

    # 4.  Tổng số Quiz yêu thích
    total_favorites = db.query(models.Quiz).filter(
        models.Quiz.owner_id == current_user.id, 
        models.Quiz.is_favorite == True
    ).count()

    return {
        "total_quizzes": total_quizzes,
        "total_completed": total_completed,
        "total_folders": total_folders,
        "total_favorites": total_favorites
    }

# API 5 QUIZ TẠO GẦN NHẤT (Dùng cho bảng 'Lịch sử tạo đề') 
@app.get("/recent-quizzes", response_model=List[schemas.QuizOut])
def get_recent_quizzes(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    recent_quizzes = db.query(models.Quiz)\
        .filter(models.Quiz.owner_id == current_user.id)\
        .order_by(models.Quiz.created_at.desc())\
        .limit(5)\
        .all()
    return recent_quizzes

#  API 5 LỊCH SỬ LÀM BÀI GẦN NHẤT
@app.get("/quiz-history")
def get_quiz_history(
    limit: int = 5,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    results = db.query(models.QuizResult)\
        .filter(models.QuizResult.user_id == current_user.id)\
        .order_by(desc(models.QuizResult.completed_at))\
        .limit(limit).all()
    
    history_list = []
    for r in results:
        history_list.append({
            "quiz_id": r.quiz_id,
            "quiz_title": r.quiz.title if r.quiz else "Quiz đã xóa",
            "score": r.score,
            "total_questions": r.total_questions,
            "completed_at": r.completed_at
        })
    return history_list

#  API DANH SÁCH YÊU THÍCH 
@app.get("/favorite-quizzes", response_model=List[schemas.QuizOut])
def get_favorite_quizzes(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    return db.query(models.Quiz).options(
        joinedload(models.Quiz.questions) 
    ).filter(
        models.Quiz.owner_id == current_user.id,
        models.Quiz.is_favorite == True
    ).all()

@app.patch("/quizzes/{quiz_id}/favorite")
def toggle_favorite(
    quiz_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id, models.Quiz.owner_id == current_user.id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz không tồn tại")
    
    quiz.is_favorite = not quiz.is_favorite 
    db.commit()
    return {"id": quiz.id, "is_favorite": quiz.is_favorite}

# API QUẢN LÝ FOLDER 
@app.post("/folders", response_model=schemas.FolderOut)
def create_folder(
    folder: schemas.FolderCreate,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    new_folder = models.Folder(name=folder.name, user_id=current_user.id)
    db.add(new_folder)
    db.commit()
    db.refresh(new_folder)
    return new_folder

@app.get("/folders", response_model=List[schemas.FolderOut])
def get_user_folders(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    return db.query(models.Folder).filter(models.Folder.user_id == current_user.id).all()

@app.delete("/folders/{folder_id}")
def delete_folder(
    folder_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    folder = db.query(models.Folder).filter(models.Folder.id == folder_id).first()
    if not folder or folder.user_id != current_user.id:
        raise HTTPException(status_code=404, detail="Không tìm thấy thư mục")
    
    quizzes_in_folder = db.query(models.Quiz).filter(models.Quiz.folder_id == folder_id).all()
    for q in quizzes_in_folder:
        q.folder_id = None
    
    db.delete(folder)
    db.commit()
    return {"message": "Đã xóa thư mục (các quiz đã được chuyển ra ngoài)"}

@app.put("/quizzes/{quiz_id}/move")
def move_quiz(
    quiz_id: int,
    move_data: schemas.MoveQuizSchema,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz or quiz.owner_id != current_user.id:
        raise HTTPException(status_code=404, detail="Không tìm thấy quiz")
    
    quiz.folder_id = move_data.folder_id
    db.commit()
    return {"message": "Đã di chuyển quiz thành công"}


#  API XUẤT QUIZ RA FILE WORD (.DOCX) - download quiz
@app.get("/quizzes/{quiz_id}/export/docx")
def export_quiz_docx(
    quiz_id: int,
    db: Session = Depends(get_db)
):
    quiz = db.query(models.Quiz).filter(models.Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Không tìm thấy quiz")

    doc = docx.Document()
    
    heading = doc.add_heading(quiz.title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph(f"Số lượng câu hỏi: {len(quiz.questions)}")
    doc.add_paragraph("-" * 50)

    for i, q in enumerate(quiz.questions):
        p = doc.add_paragraph()
        run = p.add_run(f"Câu {i+1}: {q.question_text}")
        run.bold = True
        run.font.size = Pt(12)
        
        doc.add_paragraph(f"A. {q.choice_a}")
        doc.add_paragraph(f"B. {q.choice_b}")
        doc.add_paragraph(f"C. {q.choice_c}")
        doc.add_paragraph(f"D. {q.choice_d}")
        
        doc.add_paragraph() 

    doc.add_page_break()
    doc.add_heading("ĐÁP ÁN & GIẢI THÍCH", 1)
    
    for i, q in enumerate(quiz.questions):
        p = doc.add_paragraph()
        runner = p.add_run(f"Câu {i+1}: Đáp án {q.correct_answer}")
        runner.bold = True
        
        if q.explanation:
            doc.add_paragraph(f"Giải thích: {q.explanation}")
        
        doc.add_paragraph("-" * 20)

    byte_io = io.BytesIO()
    doc.save(byte_io)
    byte_io.seek(0) 

    filename = f"Quiz_{quiz.id}.docx"
    from urllib.parse import quote
    encoded_filename = quote(filename)

    return StreamingResponse(
        byte_io,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename*=utf-8''{encoded_filename}"}
    )

#   API BIỂU ĐỒ HOẠT ĐỘNG  
@app.get("/activity-chart")
def get_activity_chart_data(
    time_range: str = "week",
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    today = date.today()
    dates = []
    quizzes_created_data = []
    quizzes_taken_data = [] 

    days_to_fetch = 30 if time_range == "month" else 7

    for i in range(days_to_fetch - 1, -1, -1):
        target_date = today - timedelta(days=i)
        dates.append(target_date.strftime("%d/%m"))

        # Đếm số Quiz tạo
        quiz_count = db.query(models.Quiz).filter(
            models.Quiz.owner_id == current_user.id,
            func.date(models.Quiz.created_at) == target_date
        ).count()
        quizzes_created_data.append(quiz_count)

        #  Đếm số Bài đã làm (QuizResult)
        taken_count = db.query(models.QuizResult).filter(
            models.QuizResult.user_id == current_user.id,
            func.date(models.QuizResult.completed_at) == target_date
        ).count()
        quizzes_taken_data.append(taken_count)

    return {
        "labels": dates,
        "created": quizzes_created_data,
        "taken": quizzes_taken_data #  Trả về key mới
    }